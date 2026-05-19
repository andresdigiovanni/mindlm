from pathlib import Path

import fitz
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation

from mindlm.core.exceptions import ParseError
from mindlm.core.models import ParsedDocument
from mindlm.core.parsing.base import DocumentParser


class RawParser(DocumentParser):
    def parse(self, path: Path) -> ParsedDocument:
        suffix = path.suffix.lower()
        match suffix:
            case ".pdf":
                doc = fitz.open(str(path))
                page_texts = [str(page.get_text()) for page in doc]
                pos = 0
                page_breaks: list[int] = []
                for i, pt in enumerate(page_texts):
                    pos += len(pt)
                    page_breaks.append(pos)
                    if i < len(page_texts) - 1:
                        pos += 1  # advance past the "\n" separator
                full_text = "\n".join(page_texts)
                return ParsedDocument(text=full_text, page_breaks=page_breaks)
            case ".html" | ".htm":
                soup = BeautifulSoup(path.read_bytes(), "lxml")
                return ParsedDocument(
                    text=str(soup.get_text(separator="\n", strip=True)),
                    page_breaks=[],
                )
            case ".md" | ".markdown" | ".txt":
                return ParsedDocument(
                    text=path.read_text(encoding="utf-8"),
                    page_breaks=[],
                )
            case ".docx":
                doc_x = DocxDocument(str(path))
                return ParsedDocument(
                    text="\n".join(para.text for para in doc_x.paragraphs),
                    page_breaks=[],
                )
            case ".pptx":
                prs = Presentation(str(path))
                texts: list[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                return ParsedDocument(text="\n".join(texts), page_breaks=[])
            case _:
                raise ParseError(str(path), "raw")
