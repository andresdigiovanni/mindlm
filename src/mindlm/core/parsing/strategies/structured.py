from pathlib import Path

import fitz
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation

from mindlm.core.exceptions import ParseError
from mindlm.core.models import ParsedDocument
from mindlm.core.parsing.base import DocumentParser


class StructuredParser(DocumentParser):
    def parse(self, path: Path) -> ParsedDocument:
        suffix = path.suffix.lower()
        match suffix:
            case ".pdf":
                doc = fitz.open(str(path))
                page_texts: list[str] = []
                for page in doc:
                    blocks = page.get_text("blocks")
                    sorted_blocks = sorted(blocks, key=lambda b: (b[1], b[0]))
                    page_text = "\n".join(
                        b[4].strip() for b in sorted_blocks if b[4].strip()
                    )
                    page_texts.append(page_text)
                pos = 0
                page_breaks: list[int] = []
                for i, pt in enumerate(page_texts):
                    pos += len(pt)
                    page_breaks.append(pos)
                    if i < len(page_texts) - 1:
                        pos += 1  # advance past the "\n" separator
                return ParsedDocument(
                    text="\n".join(page_texts), page_breaks=page_breaks
                )
            case ".html" | ".htm":
                soup = BeautifulSoup(path.read_bytes(), "lxml")
                parts: list[str] = []
                for tag in soup.find_all(
                    ["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "td", "th"]
                ):
                    name = tag.name
                    text = tag.get_text(strip=True)
                    if not text:
                        continue
                    match name:
                        case "h1":
                            parts.append(f"# {text}")
                        case "h2":
                            parts.append(f"## {text}")
                        case "h3" | "h4" | "h5" | "h6":
                            parts.append(f"### {text}")
                        case "li":
                            parts.append(f"- {text}")
                        case _:
                            parts.append(text)
                return ParsedDocument(text="\n".join(parts), page_breaks=[])
            case ".md" | ".markdown":
                return ParsedDocument(
                    text=path.read_text(encoding="utf-8"), page_breaks=[]
                )
            case ".docx":
                doc_x = DocxDocument(str(path))
                parts_d: list[str] = []
                for para in doc_x.paragraphs:
                    if not para.text.strip():
                        continue
                    style = para.style.name if para.style else ""
                    if "Heading 1" in style:
                        parts_d.append(f"# {para.text}")
                    elif "Heading 2" in style:
                        parts_d.append(f"## {para.text}")
                    else:
                        parts_d.append(para.text)
                return ParsedDocument(text="\n".join(parts_d), page_breaks=[])
            case ".pptx":
                prs = Presentation(str(path))
                slide_texts: list[str] = []
                for i, slide in enumerate(prs.slides, 1):
                    title = ""
                    body_parts: list[str] = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            if shape.shape_type == 13:  # title placeholder
                                title = shape.text_frame.text
                            else:
                                body_parts.append(shape.text_frame.text)
                    header = f"## Slide {i}: {title}" if title else f"## Slide {i}"
                    slide_texts.append(header)
                    slide_texts.extend(body_parts)
                return ParsedDocument(text="\n".join(slide_texts), page_breaks=[])
            case _:
                raise ParseError(str(path), "structured")
